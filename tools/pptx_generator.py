"""
Générateur PowerPoint - python-pptx.
Format 16:9, fond foncé, texte blanc centré, découpage par volume (50 mots max).
"""

import re
from pathlib import Path
from typing import Optional

import yaml
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

PROJECT_ROOT = Path(__file__).resolve().parent.parent
CONFIG_PATH = PROJECT_ROOT / "args" / "config.yaml"


def _load_config(config_path: Optional[Path] = None) -> dict:
    path = config_path or CONFIG_PATH
    with open(path, encoding="utf-8") as f:
        return yaml.safe_load(f)


def _strip_html(text: str) -> str:
    if not text:
        return ""
    text = re.sub(r"<[^>]+>", " ", text)
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def _word_count(text: str) -> int:
    """Compte le nombre de mots (séparés par des espaces)."""
    return len(text.split()) if text else 0


def _split_text_by_words(
    text: str,
    max_words: int = 50,
    separators: Optional[list[str]] = None,
) -> list[str]:
    """
    Découpe le texte en blocs de max_words mots maximum.
    - Priorise : point (.), point-virgule (;), virgule (,)
    - Ne coupe jamais un mot
    - Si dépassement sans ponctuation : coupe à la fin du dernier mot complet.
    """
    if not text or not text.strip():
        return []

    text = text.strip()
    config = _load_config()
    slicing = config.get("slicing", {})
    max_words = slicing.get("max_words_per_slide", 50)
    seps = separators or slicing.get("separators_priority", [". ", "; ", ", ", " "])

    # Normaliser les sauts de ligne
    text = re.sub(r"[\r\n]+", " ", text)
    text = re.sub(r"\s+", " ", text).strip()

    chunks: list[str] = []
    remaining = text

    while remaining.strip():
        remaining = remaining.strip()
        words = remaining.split()
        if len(words) <= max_words:
            chunks.append(remaining)
            break

        # Fenêtre = les max_words premiers mots (coupure avant le mot max_words+1)
        window = " ".join(words[:max_words])
        cut_pos = -1

        # Priorité : point > point-virgule > virgule > espace (coupure à la limite des mots)
        for sep in seps:
            if sep == " ":
                cut_pos = len(window)
                break
            idx = window.rfind(sep)
            if idx >= 0:
                cut_pos = idx + len(sep)
                break

        if cut_pos <= 0:
            cut_pos = len(window)

        # Ne jamais découper avant une ouverture ou fermeture de guillemets "
        rest = remaining[cut_pos:]
        next_part = rest.lstrip()
        if next_part.startswith('"'):
            close_idx = next_part.find('"', 1)
            if close_idx >= 0:
                # Inclure l'ouverture, le contenu et la fermeture dans le chunk actuel
                chars_to_add = (len(rest) - len(next_part)) + close_idx + 1
                cut_pos += chars_to_add
            else:
                # Guillemet non fermé : étendre jusqu'au prochain séparateur
                for sep in seps:
                    if sep != " ":
                        idx = next_part.find(sep, 1)
                        if idx >= 0:
                            cut_pos += (len(rest) - len(next_part)) + idx + len(sep)
                            break
                else:
                    cut_pos = len(remaining)

        chunk = remaining[:cut_pos].strip()
        remaining = remaining[cut_pos:].strip()

        if chunk:
            chunks.append(chunk)

    return chunks


def _get_slide_dimensions(config: dict) -> tuple[float, float]:
    """Retourne (width, height) en inches selon aspect_ratio."""
    pres = config.get("presentation", {})
    ratio = pres.get("aspect_ratio", "16:9")
    if ratio == "16:9":
        return 13.333, 7.5
    return 10.0, 7.5  # 4:3 par défaut


def _get_theme_colors(config: dict, theme: str) -> tuple[str, str, str, str]:
    """Retourne (bg_lecture, bg_chant, text_color, title_color) pour le thème."""
    design = config.get("design", {})
    theme_colors = design.get("theme_colors", {}).get(theme)
    if theme_colors:
        bg = theme_colors.get("background", {})
        return (
            bg.get("color", "#0a1628"),
            bg.get("color_chant", "#0d2818"),
            theme_colors.get("text", "#FFFFFF"),
            theme_colors.get("title_rappel", "#CCCCCC"),
        )
    bg = design.get("background", {})
    text_cfg = design.get("text", {})
    title_cfg = design.get("title_rappel", {})
    return (
        bg.get("color", "#0a1628"),
        bg.get("color_chant", "#0d2818"),
        text_cfg.get("color", "#FFFFFF"),
        title_cfg.get("color", "#CCCCCC"),
    )


def _add_slide(
    prs: Presentation,
    config: dict,
    title: str,
    body: str,
    is_continuation: bool = False,
    slide_type: str = "lecture",
    theme: str = "dark",
) -> None:
    """Ajoute une slide 16:9 avec texte centré. slide_type: 'lecture' | 'chant' | 'message'. theme: 'dark' | 'light'."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    bg_lecture, bg_chant, font_color_hex, title_color_hex = _get_theme_colors(config, theme)
    if slide_type == "chant":
        bg_color = bg_chant
    else:
        bg_color = bg_lecture
    r = int(bg_color[1:3], 16)
    g = int(bg_color[3:5], 16)
    b = int(bg_color[5:7], 16)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)

    # Style texte (couleurs selon le thème)
    text_cfg = config.get("design", {}).get("text", {})
    font_name = text_cfg.get("font", "Calibri")
    font_size = text_cfg.get("size", 34)
    tr = int(font_color_hex[1:3], 16)
    tg = int(font_color_hex[3:5], 16)
    tb = int(font_color_hex[5:7], 16)

    title_cfg = config.get("design", {}).get("title_rappel", {})
    title_size = title_cfg.get("size", 20)
    tr2 = int(title_color_hex[1:3], 16)
    tg2 = int(title_color_hex[3:5], 16)
    tb2 = int(title_color_hex[5:7], 16)

    width_in, height_in = _get_slide_dimensions(config)
    margin = 0.6

    # Zone utilisable
    left = Inches(margin)
    top = Inches(margin)
    width = Inches(width_in - 2 * margin)
    height = Inches(height_in - 2 * margin)

    # Titre de rappel (centré)
    title_height = Inches(0.7)
    tf = slide.shapes.add_textbox(left, top, width, title_height)
    tf.text_frame.word_wrap = True
    p = tf.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(title_size)
    p.font.name = font_name
    p.font.color.rgb = RGBColor(tr2, tg2, tb2)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Corps centré (horizontal et vertical)
    body_top = Inches(margin + 0.8)
    body_height = Inches(height_in - 2 * margin - 1.0)
    body_box = slide.shapes.add_textbox(left, body_top, width, body_height)
    body_box.text_frame.word_wrap = True
    body_box.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = body_box.text_frame.paragraphs[0]
    p.text = body
    p.font.size = Pt(font_size)
    p.font.name = font_name
    p.font.color.rgb = RGBColor(tr, tg, tb)
    p.alignment = PP_ALIGN.CENTER


def generate_pptx(
    blocs: list[dict],
    output_path: Path,
    config_path: Optional[Path] = None,
    theme: str = "dark",
) -> Path:
    """
    Génère un fichier PowerPoint à partir d'une liste de blocs.
    theme: 'dark' (fond foncé, texte clair) ou 'light' (fond clair, texte foncé).
    """
    cfg_path = config_path or CONFIG_PATH
    config = _load_config(cfg_path)

    width_in, height_in = _get_slide_dimensions(config)
    prs = Presentation()
    prs.slide_width = Inches(width_in)
    prs.slide_height = Inches(height_in)

    total_slides = 0

    for bloc in blocs:
        t = bloc.get("type", "")
        label = ""

        if t == "lecture":
            ref = bloc.get("reference", "")
            intro = bloc.get("intro_lue", "")
            contenu = _strip_html(bloc.get("contenu", ""))
            titre_base = intro or ref
            label = titre_base or "Lecture"
            chunks = _split_text_by_words(contenu)

        elif t == "chant":
            titre = bloc.get("titre", "Chant")
            paroles = bloc.get("paroles", "")
            label = titre
            chunks = _split_text_by_words(paroles)

        elif t == "message":
            titre = bloc.get("titre", "Message")
            contenu = bloc.get("contenu", "")
            label = titre
            chunks = _split_text_by_words(contenu)

        else:
            continue

        n = len(chunks)
        total_slides += n
        print(f"[pptx] {label} : {n} slide(s) générée(s)", flush=True)

        for i, chunk in enumerate(chunks):
            titre_slide = f"{label} (suite)" if i > 0 else label
            if t == "lecture":
                _add_slide(prs, config, titre_slide, chunk, is_continuation=(i > 0), slide_type="lecture", theme=theme)
            elif t == "chant":
                _add_slide(prs, config, titre_slide, chunk, is_continuation=(i > 0), slide_type="chant", theme=theme)
            else:
                _add_slide(prs, config, titre_slide, chunk, is_continuation=(i > 0), slide_type="message", theme=theme)

    print(f"[pptx] Total : {total_slides} slide(s)", flush=True)

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path
